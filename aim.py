from pptx import Presentation
from pptx.util import Inches, Pt

# Create presentation
prs = Presentation()

# Define slide layout
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]
blank_slide_layout = prs.slide_layouts[6]

# Helper function to add notes
def add_notes(slide, text):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = text

# Slide 1 – Titel & Introductie
slide1 = prs.slides.add_slide(title_slide_layout)
slide1.shapes.title.text = "Slimme Asset & Inventory Management met Discovery-integratie"
slide1.placeholders[1].text = "Een toekomstgerichte architectuur voor datagedreven processen"
add_notes(slide1, "Welkom bij deze presentatie over een slimme aanpak voor Asset & Inventory Management met Discovery-integratie. We bekijken hoe deze architectuur processen kan verbeteren.")

# Slide 2 – Uitdagingen in Asset Management
slide2 = prs.slides.add_slide(bullet_slide_layout)
slide2.shapes.title.text = "Uitdagingen in Asset Management"
content = slide2.shapes.placeholders[1].text_frame
for point in ["Versnipperde data", "Verouderde inventarisaties", "Geen inzicht in relaties", "Risico’s voor compliance"]:
    content.add_paragraph().text = point
add_notes(slide2, "Veel organisaties kampen met versnipperde data en verouderde inventarisaties. Dit belemmert inzicht en verhoogt compliance-risico’s.")

# Slide 3 – Architectuuroverzicht
slide3 = prs.slides.add_slide(blank_slide_layout)
title_shape = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
title_frame = title_shape.text_frame
p = title_frame.add_paragraph()
p.text = "Architectuuroverzicht"
p.font.size = Pt(32)
p.font.bold = True

text_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1))
frame = text_box.text_frame
frame.text = "Geïntegreerd systeem dat data verzamelt, verifieert en distribueert"
add_notes(slide3, "Deze architectuur verbindt verschillende interfaces met een centrale database en een Discovery-engine, die data verzamelt en analyseert.")

# Slide 4 – Discovery Engine
slide4 = prs.slides.add_slide(bullet_slide_layout)
slide4.shapes.title.text = "Discovery Engine"
content = slide4.shapes.placeholders[1].text_frame
for point in ["Netwerk scan", "Objectherkenning", "Relatie-analyse", "Dataverificatie"]:
    content.add_paragraph().text = point
add_notes(slide4, "De Discovery-engine voert netwerk scans uit, herkent objecten, analyseert relaties en verifieert data automatisch.")

# Slide 5 – Procesintegratie
slide5 = prs.slides.add_slide(blank_slide_layout)
title_shape = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
title_frame = title_shape.text_frame
p = title_frame.add_paragraph()
p.text = "Procesintegratie"
p.font.size = Pt(32)
p.font.bold = True

text_box = slide5.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1))
frame = text_box.text_frame
frame.text = "Tabel met procesgebieden en datagebruik"
add_notes(slide5, "Data uit de centrale database stroomt naar verschillende procesgebieden zoals Finance, Operations en Support, wat integratie bevordert.")

# Slide 6 – Voordelen van het Model
slide6 = prs.slides.add_slide(bullet_slide_layout)
slide6.shapes.title.text = "Voordelen van het Model"
content = slide6.shapes.placeholders[1].text_frame
for point in ["Centrale bron van waarheid", "Automatische verificatie", "Betere besluitvorming", "Schaalbaarheid"]:
    content.add_paragraph().text = point
add_notes(slide6, "Dit model biedt een centrale bron van waarheid, automatische verificatie en schaalbaarheid voor betere besluitvorming.")

# Slide 7 – Volgende Stappen
slide7 = prs.slides.add_slide(bullet_slide_layout)
slide7.shapes.title.text = "Volgende Stappen"
content = slide7.shapes.placeholders[1].text_frame
for point in ["Proof of Concept", "Integratie", "Training", "Monitoring"]:
    content.add_paragraph().text = point
add_notes(slide7, "De volgende stappen zijn het opzetten van een Proof of Concept, integratie in bestaande systemen, training van gebruikers en continue monitoring.")

# Save presentation
prs.save("/mnt/data/Asset_Inventory_Discovery_Presentation.pptx")